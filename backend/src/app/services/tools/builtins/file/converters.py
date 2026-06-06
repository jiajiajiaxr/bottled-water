from __future__ import annotations

import csv
import hashlib
import html
import io
import json
import re
import shutil
import subprocess
import zipfile
from dataclasses import dataclass
from html.parser import HTMLParser
from pathlib import Path
from typing import Any

from app.core.errors import ValidationAppError
from app.services.document_model import normalize_document_model, render_docx, render_pdf


TEXT_EXTENSIONS = {
    ".txt",
    ".md",
    ".markdown",
    ".py",
    ".ts",
    ".tsx",
    ".js",
    ".jsx",
    ".json",
    ".html",
    ".htm",
    ".css",
    ".xml",
    ".csv",
    ".yml",
    ".yaml",
}

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".svg"}

TEXT_TYPES = {
    "text/plain",
    "text/markdown",
    "application/json",
    "application/xml",
    "text/html",
    "text/css",
    "text/javascript",
    "application/javascript",
    "text/csv",
}


@dataclass(frozen=True)
class GeneratedFile:
    content: bytes
    media_type: str
    filename: str


class _HTMLTextParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []

    def handle_starttag(self, tag: str, _attrs: list[tuple[str, str | None]]) -> None:
        if tag in {"p", "div", "section", "article", "li", "tr", "h1", "h2", "h3", "br"}:
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        value = data.strip()
        if value:
            self.parts.append(value)

    def text(self) -> str:
        text = " ".join(self.parts)
        text = re.sub(r"\s*\n\s*", "\n", text)
        text = re.sub(r"[ \t]{2,}", " ", text)
        return text.strip()


def _decode(raw: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "latin-1"):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="ignore")


def _text_from_html(raw: bytes) -> str:
    parser = _HTMLTextParser()
    parser.feed(_decode(raw))
    return parser.text()


def _text_from_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except Exception as exc:  # pragma: no cover - optional dependency diagnostics.
        raise ValidationAppError("PDF 解析依赖 pypdf 不可用") from exc
    reader = PdfReader(str(path))
    parts = []
    for index, page in enumerate(reader.pages[:80], start=1):
        text = page.extract_text() or ""
        if text.strip():
            parts.append(f"--- page {index} ---\n{text}")
    return "\n\n".join(parts)


def _text_from_docx(path: Path) -> str:
    try:
        from docx import Document
    except Exception as exc:  # pragma: no cover
        raise ValidationAppError("Word 解析依赖 python-docx 不可用") from exc
    document = Document(str(path))
    paragraphs = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    table_rows: list[str] = []
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                table_rows.append(" | ".join(cells))
    return "\n".join([*paragraphs, *table_rows])


def _text_from_xlsx(path: Path) -> str:
    try:
        from openpyxl import load_workbook
    except Exception as exc:  # pragma: no cover
        raise ValidationAppError("Excel 解析依赖 openpyxl 不可用") from exc
    workbook = load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in workbook.worksheets[:8]:
        parts.append(f"# Sheet: {sheet.title}")
        for row in sheet.iter_rows(max_row=80, values_only=True):
            values = [str(cell) for cell in row if cell is not None]
            if values:
                parts.append(" | ".join(values))
    workbook.close()
    return "\n".join(parts)


def _text_from_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except Exception as exc:  # pragma: no cover
        raise ValidationAppError("PPT 解析依赖 python-pptx 不可用") from exc
    deck = Presentation(str(path))
    parts: list[str] = []
    for index, slide in enumerate(deck.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            parts.append(f"# Slide {index}\n" + "\n".join(texts))
    return "\n\n".join(parts)


def _text_from_image(path: Path) -> tuple[str, dict[str, Any]]:
    """Extract text from an image through a local OCR binary when available.

    The platform keeps image understanding honest: if OCR is not installed we
    return a clear degraded status instead of pretending the model saw pixels.
    """

    binary = shutil.which("tesseract")
    if not binary:
        return "", {
            "extractor": "image_ocr",
            "vision_status": "missing_tesseract",
            "ocr_available": False,
            "setup_hint": "安装 Tesseract OCR 后可自动识别 PNG/JPG 截图文字。",
        }
    try:
        completed = subprocess.run(
            [binary, str(path), "stdout", "-l", "chi_sim+eng"],
            capture_output=True,
            text=True,
            encoding="utf-8",
            errors="ignore",
            timeout=20,
            check=False,
        )
    except subprocess.TimeoutExpired:
        return "", {
            "extractor": "image_ocr",
            "vision_status": "timeout",
            "ocr_available": True,
            "error": "OCR 识别超时",
        }
    except Exception as exc:  # pragma: no cover - local binary diagnostics.
        return "", {
            "extractor": "image_ocr",
            "vision_status": "failed",
            "ocr_available": True,
            "error": str(exc),
        }

    text = (completed.stdout or "").strip()
    metadata: dict[str, Any] = {
        "extractor": "image_ocr",
        "vision_status": "parsed" if text else "empty",
        "ocr_available": True,
        "exit_code": completed.returncode,
    }
    if completed.returncode != 0 and not text:
        metadata["vision_status"] = "failed"
        metadata["error"] = (completed.stderr or "").strip()[:800]
    return text, metadata


def extract_text_from_path(path: Path, *, content_type: str = "", filename: str = "") -> dict[str, Any]:
    suffix = (path.suffix or Path(filename).suffix).lower()
    raw = path.read_bytes()
    metadata: dict[str, Any] = {"extension": suffix, "extractor": "none", "byte_size": len(raw)}
    text = ""
    if content_type in TEXT_TYPES or suffix in TEXT_EXTENSIONS:
        text = _text_from_html(raw) if suffix in {".html", ".htm"} or content_type == "text/html" else _decode(raw)
        metadata["extractor"] = "native_text"
    elif suffix == ".pdf" or content_type == "application/pdf":
        text = _text_from_pdf(path)
        metadata["extractor"] = "pypdf"
    elif suffix == ".docx" or content_type.endswith("wordprocessingml.document"):
        text = _text_from_docx(path)
        metadata["extractor"] = "python-docx"
    elif suffix == ".xlsx" or content_type.endswith("spreadsheetml.sheet"):
        text = _text_from_xlsx(path)
        metadata["extractor"] = "openpyxl"
    elif suffix == ".pptx" or content_type.endswith("presentationml.presentation"):
        text = _text_from_pptx(path)
        metadata["extractor"] = "python-pptx"
    elif suffix in IMAGE_EXTENSIONS or content_type.startswith("image/"):
        text, image_metadata = _text_from_image(path)
        metadata.update(image_metadata)
    else:
        metadata["extractor"] = "unsupported"
    text = text.strip()
    return {
        "status": "parsed" if text else "stored",
        "text": text[:300_000],
        "metadata": metadata,
    }


def preview_payload(path: Path, *, content_type: str = "", filename: str = "") -> dict[str, Any]:
    extracted = extract_text_from_path(path, content_type=content_type, filename=filename)
    suffix = (path.suffix or Path(filename).suffix).lower()
    mode = "text"
    if suffix == ".pdf" or content_type == "application/pdf":
        mode = "pdf"
    elif suffix in IMAGE_EXTENSIONS or content_type.startswith("image/"):
        mode = "image"
    elif suffix in {".docx", ".xlsx", ".pptx"}:
        mode = "office_text"
    return {
        "filename": filename or path.name,
        "content_type": content_type or "application/octet-stream",
        "mode": mode,
        "preview_text": extracted["text"][:30_000],
        "metadata": extracted["metadata"],
    }


def summarize_text(text: str, *, max_chars: int = 1200) -> str:
    cleaned = re.sub(r"\s+", " ", text).strip()
    if len(cleaned) <= max_chars:
        return cleaned
    sentences = re.split(r"(?<=[。！？.!?])\s+", cleaned)
    summary = ""
    for sentence in sentences:
        if len(summary) + len(sentence) > max_chars:
            break
        summary += sentence + " "
    return (summary.strip() or cleaned[:max_chars]).strip()


def embed_text(text: str, *, dimensions: int = 32) -> list[float]:
    digest = hashlib.sha256(text.encode("utf-8", errors="ignore")).digest()
    values = []
    for index in range(dimensions):
        byte = digest[index % len(digest)]
        values.append(round((byte / 255.0) * 2 - 1, 6))
    return values


def _html_document(title: str, body: str) -> str:
    paragraphs = "\n".join(f"<p>{html.escape(line)}</p>" for line in body.splitlines() if line.strip())
    return (
        "<!doctype html><html><head><meta charset=\"utf-8\">"
        f"<title>{html.escape(title)}</title>"
        "<style>body{font-family:system-ui,sans-serif;line-height:1.7;padding:40px;color:#111827}"
        "main{max-width:860px;margin:auto}h1{font-size:32px}</style></head><body><main>"
        f"<h1>{html.escape(title)}</h1>{paragraphs}</main></body></html>"
    )


def generate_docx(title: str, body: str, content_model: dict[str, Any] | None = None) -> bytes:
    model = normalize_document_model(content_model, title=title, source_text=body)
    return render_docx(model)


def generate_xlsx(title: str, body: str) -> bytes:
    try:
        from openpyxl import Workbook
    except Exception as exc:  # pragma: no cover
        raise ValidationAppError("Excel 生成依赖 openpyxl 不可用") from exc
    buffer = io.BytesIO()
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "AgentHub"
    sheet.append([title])
    sheet.append([])
    for index, line in enumerate([line for line in body.splitlines() if line.strip()], start=1):
        sheet.append([index, line.strip()])
    sheet.column_dimensions["A"].width = 12
    sheet.column_dimensions["B"].width = 96
    workbook.save(buffer)
    workbook.close()
    return buffer.getvalue()


def generate_pptx(title: str, body: str) -> bytes:
    try:
        from pptx import Presentation
    except Exception as exc:  # pragma: no cover
        raise ValidationAppError("PPT 生成依赖 python-pptx 不可用") from exc
    buffer = io.BytesIO()
    deck = Presentation()
    title_slide = deck.slides.add_slide(deck.slide_layouts[0])
    title_slide.shapes.title.text = title
    title_slide.placeholders[1].text = "Generated by AgentHub"
    lines = [line.strip() for line in body.splitlines() if line.strip()]
    for start in range(0, max(len(lines), 1), 5):
        slide = deck.slides.add_slide(deck.slide_layouts[1])
        slide.shapes.title.text = f"{title} · {start // 5 + 1}"
        body_shape = slide.placeholders[1]
        body_shape.text = "\n".join(lines[start : start + 5]) or body[:300]
    deck.save(buffer)
    return buffer.getvalue()


def generate_pdf(title: str, body: str, content_model: dict[str, Any] | None = None) -> bytes:
    model = normalize_document_model(content_model, title=title, source_text=body)
    return render_pdf(model)


def generate_markdown(title: str, body: str) -> bytes:
    return f"# {title}\n\n{body.strip()}\n".encode("utf-8")


def generate_html(title: str, body: str) -> bytes:
    return _html_document(title, body).encode("utf-8")


def generate_zip(filename: str, files: dict[str, str | bytes]) -> bytes:
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("README.md", f"# {filename}\n\nGenerated by AgentHub.\n")
        for path, value in files.items():
            archive.writestr(path, value)
    return buffer.getvalue()


GENERATORS = {
    "pdf": (generate_pdf, "application/pdf"),
    "docx": (generate_docx, "application/vnd.openxmlformats-officedocument.wordprocessingml.document"),
    "xlsx": (generate_xlsx, "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"),
    "pptx": (generate_pptx, "application/vnd.openxmlformats-officedocument.presentationml.presentation"),
    "html": (generate_html, "text/html; charset=utf-8"),
    "md": (generate_markdown, "text/markdown; charset=utf-8"),
    "markdown": (generate_markdown, "text/markdown; charset=utf-8"),
}


def generate_file(
    format_name: str,
    *,
    title: str,
    body: str,
    content_model: dict[str, Any] | None = None,
) -> GeneratedFile:
    fmt = format_name.lower().strip(".")
    if fmt not in GENERATORS:
        raise ValidationAppError(f"不支持的生成格式：{format_name}")
    generator, media_type = GENERATORS[fmt]
    normalized = "md" if fmt == "markdown" else fmt
    content = generator(title, body, content_model) if normalized in {"pdf", "docx"} else generator(title, body)
    safe = re.sub(r"[^A-Za-z0-9_.\-\u4e00-\u9fff]+", "-", title).strip("-") or "agenthub-artifact"
    return GeneratedFile(content=content, media_type=media_type, filename=f"{safe[:70]}.{normalized}")


def convert_file(path: Path, *, content_type: str, filename: str, target_format: str) -> GeneratedFile:
    extracted = extract_text_from_path(path, content_type=content_type, filename=filename)
    title = Path(filename or path.name).stem or "AgentHub 文件转换"
    body = extracted["text"] or f"{filename or path.name} 暂无可提取文本。"
    if target_format.lower().strip(".") == "json":
        content = json.dumps({"filename": filename, "content_type": content_type, "text": body}, ensure_ascii=False, indent=2).encode("utf-8")
        return GeneratedFile(content=content, media_type="application/json; charset=utf-8", filename=f"{title}.json")
    if target_format.lower().strip(".") == "csv":
        buffer = io.StringIO()
        writer = csv.writer(buffer)
        writer.writerow(["line", "content"])
        for index, line in enumerate(body.splitlines(), start=1):
            writer.writerow([index, line])
        return GeneratedFile(content=buffer.getvalue().encode("utf-8-sig"), media_type="text/csv; charset=utf-8", filename=f"{title}.csv")
    return generate_file(target_format, title=title, body=body)
